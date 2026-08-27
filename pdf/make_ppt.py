#!/usr/bin/env python3
"""Build the defense deck for 项目6 - Agent for Kernel.

Every number in this deck comes from a recorded measurement under
/tmp/apxinf-evidence/ or from REPORT.md in the submitted commit; nothing is
estimated. Layout is done with explicit geometry (no template dependency) so
the file renders identically wherever it is opened.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

import os

# 作品名：4090 能飞（改编自网络热梗「中国人能飞」）
AUTHOR = os.environ.get("PPT_AUTHOR", "程仁龙")
WORK = "4090能飞"
OUT = Path(f"/mnt/chuangxin/team2/ApxInf/pdf/项目6_{AUTHOR}_{WORK}_v1.pptx")

# 16:9
W, H = Inches(13.333), Inches(7.5)

INK = RGBColor(0x1A, 0x1D, 0x24)      # near-black text
MUTED = RGBColor(0x5A, 0x63, 0x73)    # secondary text
ACCENT = RGBColor(0x00, 0x66, 0xCC)   # primary accent
GOOD = RGBColor(0x0E, 0x8A, 0x4F)     # improvement / pass
WARN = RGBColor(0xC0, 0x39, 0x2B)     # limitation / reject
RULE = RGBColor(0xDD, 0xE1, 0xE8)     # hairlines
PANEL = RGBColor(0xF5, 0xF7, 0xFA)    # panel fill
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

CN = "SimHei"
MONO = "SimHei"

from pptx.oxml.ns import qn


def _apply_font(run, name):
    """Set the typeface for both latin and east-asian scripts, so CJK text
    actually renders in the requested font instead of the theme default."""
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, *, size=14, bold=False, color=INK, font=CN, space_after=4,
         align=PP_ALIGN.LEFT, first=False, line=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    if line:
        p.line_spacing = line
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    _apply_font(run, font)
    return p


def rect(slide, x, y, w, h, fill=PANEL, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.04
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    shape.text_frame.text = ""
    return shape


def hairline(slide, x, y, w):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(1.2))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def header(slide, title, kicker=None):
    """Standard slide header: small kicker, title, accent rule."""
    y = Inches(0.45)
    if kicker:
        tf = textbox(slide, Inches(0.7), y, Inches(12), Inches(0.28))
        para(tf, kicker, size=12, bold=True, color=ACCENT, first=True)
        y = Inches(0.75)
    tf = textbox(slide, Inches(0.7), y, Inches(12), Inches(0.55))
    para(tf, title, size=27, bold=True, first=True)
    hairline(slide, Inches(0.7), y + Inches(0.62), Inches(1.5))
    return y + Inches(0.95)


def table(slide, x, y, w, rows, col_w, *, head_size=12, body_size=12,
          row_h=Inches(0.34), colors=None):
    """Lightweight table: first row is the header. `colors` maps row index to
    a per-row text color override."""
    cur = y
    for r, row in enumerate(rows):
        is_head = r == 0
        if is_head:
            bg = rect(slide, x, cur, w, row_h, fill=RGBColor(0xEC, 0xF1, 0xF8))
        elif r % 2 == 0:
            bg = rect(slide, x, cur, w, row_h, fill=RGBColor(0xFA, 0xFB, 0xFD))
        cx = x
        for c, cell in enumerate(row):
            cw = col_w[c]
            tf = textbox(slide, cx + Inches(0.09), cur, cw - Inches(0.18), row_h,
                         anchor=MSO_ANCHOR.MIDDLE)
            color = INK if is_head else (colors or {}).get(r, INK)
            align = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
            para(tf, str(cell), size=head_size if is_head else body_size,
                 bold=is_head, color=color, align=align, first=True,
                 font=MONO if (not is_head and c > 0) else CN)
            cx += cw
        cur += row_h
    return cur


def metric(slide, x, y, w, big, label, sub=None, color=ACCENT):
    h = Inches(1.5) if sub else Inches(1.25)
    rect(slide, x, y, w, h, fill=WHITE, line=RULE)
    tf = textbox(slide, x + Inches(0.16), y + Inches(0.14), w - Inches(0.32), Inches(0.62))
    para(tf, big, size=34, bold=True, color=color, first=True, align=PP_ALIGN.CENTER)
    tf = textbox(slide, x + Inches(0.16), y + Inches(0.78), w - Inches(0.32), Inches(0.28))
    para(tf, label, size=12, bold=True, color=INK, first=True, align=PP_ALIGN.CENTER)
    if sub:
        tf = textbox(slide, x + Inches(0.16), y + Inches(1.06), w - Inches(0.32), Inches(0.34))
        para(tf, sub, size=10, color=MUTED, first=True, align=PP_ALIGN.CENTER)


def bullets(slide, x, y, w, items, *, size=14, gap=9, bullet="—"):
    # Reserve only the space that actually remains above the footer, so the
    # placeholder box never extends past the canvas.
    avail = H - y - Inches(0.62)
    tf = textbox(slide, x, y, w, max(Inches(0.4), avail))
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            lead, rest = item
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(gap)
            p.line_spacing = 1.25
            r1 = p.add_run()
            r1.text = f"{bullet} {lead}"
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.color.rgb = INK
            _apply_font(r1, CN)
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(size)
            r2.font.color.rgb = MUTED
            _apply_font(r2, CN)
        else:
            para(tf, f"{bullet} {item}", size=size, color=MUTED, first=(i == 0),
                 space_after=gap, line=1.25)
    return tf


def footer(slide, text):
    tf = textbox(slide, Inches(0.7), Inches(6.95), Inches(12), Inches(0.3))
    para(tf, text, size=9.5, color=MUTED, first=True)


def build() -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    # ───────────────────────── 1. Cover ─────────────────────────
    s = blank(prs)
    rect(s, 0, 0, W, Inches(2.55), fill=RGBColor(0x0F, 0x2A, 0x4A))
    tf = textbox(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.4))
    para(tf, "项目6 · Agent for Kernel", size=15, bold=True,
         color=RGBColor(0x8F, 0xC2, 0xF0), first=True)
    tf = textbox(s, Inches(0.9), Inches(0.96), Inches(11.5), Inches(0.78))
    para(tf, "4090 能飞", size=44, bold=True, color=WHITE, first=True)
    tf = textbox(s, Inches(0.9), Inches(1.72), Inches(11.5), Inches(0.36))
    para(tf, "Qwen3.8-27B W4A16 单卡 RTX 4090 推理引擎", size=19, bold=True,
         color=RGBColor(0xC5, 0xD9, 0xEE), first=True)
    tf = textbox(s, Inches(0.9), Inches(2.08), Inches(11.5), Inches(0.36))
    para(tf, "一张卡、4.6 GiB 余量，从零实现 64 层混合架构；16K 首 token 从 27 分钟到 34 秒",
         size=13.5, color=RGBColor(0x9F, 0xBE, 0xDE), first=True)

    y = Inches(2.95)
    for i, (k, v) in enumerate([
        ("架构", "48 GDN + 16 全注意力，共 64 层"),
        ("量化", "compressed-tensors W4A16 group-32 asymmetric"),
        ("硬件", "单张 RTX 4090 (SM89, 24 GiB)，TP=PP=DP=1"),
        ("提交", "commit 06993a2d2642c6f7177b57493b797d5d537e4d64"),
    ]):
        tf = textbox(s, Inches(0.9) + Inches(3.05) * (i % 2),
                     y + Inches(0.62) * (i // 2), Inches(3.0), Inches(0.55))
        para(tf, k, size=11, bold=True, color=ACCENT, first=True, space_after=1)
        para(tf, v, size=11.5, color=MUTED,
             font=MONO if k == "提交" else CN, space_after=0)
    tf = textbox(s, Inches(7.0), y, Inches(5.0), Inches(0.55))
    para(tf, "团队", size=11, bold=True, color=ACCENT, first=True, space_after=1)
    para(tf, "程仁龙（负责人）· 王天民", size=11.5, color=MUTED, space_after=0)
    tf = textbox(s, Inches(7.0), y + Inches(0.62), Inches(5.3), Inches(0.55))
    para(tf, "命名说明", size=11, bold=True, color=ACCENT, first=True, space_after=1)
    para(tf, "商品名 Qwen3.8；checkpoint 内 HF 架构为 qwen3_5，代码模块因此名为 qwen35",
         size=11.5, color=MUTED, space_after=0)

    rect(s, Inches(0.9), Inches(4.42), Inches(11.5), Inches(1.86), fill=PANEL)
    tf = textbox(s, Inches(1.15), Inches(4.6), Inches(11.0), Inches(1.5))
    para(tf, "汇报主线", size=13, bold=True, first=True, space_after=6)
    para(tf, "① 取得文本 eligibility：协议 12/12、公开功能 6/6、200 请求 soak 100%、7 个性能 cell 全部有效",
         size=12.5, color=MUTED, space_after=4)
    para(tf, "② 用可回滚的配对实验做性能工程：九项接受、八项拒绝，TTFT 约 47 倍、TPOT 约 2 倍",
         size=12.5, color=MUTED, space_after=4)
    para(tf, "③ 交付双 bonus：多模态图文链路（探针 4/4、文本数值零变化）+ C4 四路并发"
             "（官方校准全门通过，goodput 23.5 tok/s）",
         size=12.5, color=MUTED, space_after=4)
    para(tf, "④ 评测边界：无平台批准的 scorer 产物，不声称任何评分结论",
         size=12.5, color=WARN, space_after=0)
    footer(s, "所有数据均来自实测记录（随附证据档案与 REPORT.md），无估算值")

    # ───────────────────────── 2. Results at a glance ─────────────────────────
    s = blank(prs)
    y = header(s, "成果总览", "RESULTS")
    metric(s, Inches(0.7), y, Inches(2.85), "47×", "TTFT 提升",
           "1K：78.09 s → 1.63 s", GOOD)
    metric(s, Inches(3.75), y, Inches(2.85), "2.0×", "TPOT 提升",
           "1K：133.7 ms → 66.5 ms", GOOD)
    metric(s, Inches(6.80), y, Inches(2.85), "7/7", "性能 cell 有效",
           "success 5/5，最差 CV 8.1%", ACCENT)
    metric(s, Inches(9.85), y, Inches(2.85), "100%", "200 请求 soak",
           "五项 reliability 全为真", ACCENT)

    y2 = y + Inches(1.75)
    tf = textbox(s, Inches(0.7), y2, Inches(5.6), Inches(0.3))
    para(tf, "七个基础 cell：优化前 → 优化后", size=13, bold=True, first=True)
    rows = [
        ("cell", "优化前", "优化后", "倍数"),
        ("TTFT 1K", "78.09 s", "1.63 s", "47.9×"),
        ("TTFT 2K", "159.19 s", "3.31 s", "48.1×"),
        ("TTFT 4K", "320.40 s", "7.01 s", "45.7×"),
        ("TTFT 8K", "676.52 s", "14.81 s", "45.7×"),
        ("TTFT 16K", "1635.15 s", "34.58 s", "47.3×"),
        ("TPOT 1K", "133.7 ms", "66.5 ms", "2.01×"),
        ("TPOT 8K", "153.4 ms", "67.3 ms", "2.28×"),
    ]
    cw = [Inches(1.5), Inches(1.4), Inches(1.3), Inches(1.1)]
    table(s, Inches(0.7), y2 + Inches(0.38), Inches(5.3), rows, cw, row_h=Inches(0.315))

    tf = textbox(s, Inches(6.5), y2, Inches(6.1), Inches(0.3))
    para(tf, "correctness 与稳定性门禁", size=13, bold=True, first=True)
    bullets(s, Inches(6.5), y2 + Inches(0.4), Inches(6.0), [
        ("协议 gate 12/12 ", "7 个负控全部 HTTP 400 + JSON error"),
        ("公开功能 6/6 精确 ", "含三个 8K longdoc，EOS 提前终止正常"),
        ("200 请求混合 soak 100% ", "无 OOM / NaN / fallback / Xid"),
        ("proxy hidden 11/12 ", "达到资格线（非官方代理集，已标注）"),
        ("多模态 + C4 双 bonus ", "图文探针 4/4；C4 官方校准全门通过"),
        ("峰值显存 ", "文本 19958 / 多模态 20914，上限 24564 MiB"),
    ], size=12.5, gap=10)
    footer(s, "口径遵循合同：warmup 1 次 + 测 5 次取中位数，TTFT/TPOT 的 CV 均 ≤ 10%")

    # ───────────────────────── 3. Task & constraints ─────────────────────────
    s = blank(prs)
    y = header(s, "任务边界与合同约束", "SCOPE")
    tf = textbox(s, Inches(0.7), y, Inches(11.9), Inches(0.5))
    para(tf, "starter 仓库没有 Qwen3.8、没有 W4A16、没有 GDN、也没有 HTTP/SSE —— "
             "本项目是从零建立的 vertical slice，而非在既有模块上接线。",
         size=14, color=INK, first=True, line=1.3)

    y += Inches(0.62)
    left = [
        ("64 层混合主干", "每四层三个 GDN（线性注意力）+ 一个全注意力，两类层请求状态互不相同"),
        ("GDN 层", "16 key / 48 value head、head_dim 128，因果卷积 ring buffer + FP32 递归态"),
        ("全注意力层", "24 Q / 4 KV head、head_dim 256，输出门控 + partial RoPE 仅前 64 维"),
        ("逐模块混合量化", "MLP 与投影为 packed W4，in_proj 与 conv/norm 为 BF16"),
    ]
    right = [
        ("单卡固定", "禁止其他 GPU、CPU、vLLM、Transformers 作为服务 fallback"),
        ("不改评测件", "evaluation/ 下合同、生成器、scorer 一律不动"),
        ("不硬编码", "不按 case ID、公开 token 或已知答案特判"),
        ("health 必须真实", "capability 未经验证前一律 fail closed"),
    ]
    tf = textbox(s, Inches(0.7), y, Inches(5.8), Inches(0.3))
    para(tf, "模型结构要点", size=13, bold=True, color=ACCENT, first=True)
    bullets(s, Inches(0.7), y + Inches(0.38), Inches(5.6), left, size=12.5, gap=11)
    tf = textbox(s, Inches(6.8), y, Inches(5.8), Inches(0.3))
    para(tf, "合同硬约束", size=13, bold=True, color=WARN, first=True)
    bullets(s, Inches(6.8), y + Inches(0.38), Inches(5.6), right, size=12.5, gap=11)
    footer(s, "分层原则：模型知道层序、状态与融合决策；CUDA backend 只暴露设备管理与单 kernel/单库调用")

    # ───────────────────────── 4. Layer 1 evidence ─────────────────────────
    s = blank(prs)
    y = header(s, "第一层：取得文本 eligibility", "CORRECTNESS · RELIABILITY")
    tf = textbox(s, Inches(0.7), y, Inches(11.9), Inches(0.42))
    para(tf, "eligibility 是硬门：任一 reliability boolean 为假会直接令 eligible=false，"
             "不是从 10 分里扣分。因此先把门禁全部打绿，再谈性能。",
         size=13.5, color=INK, first=True, line=1.3)

    y += Inches(0.58)
    rows = [
        ("门禁项", "要求", "实测结果"),
        ("协议 gate", "malformed 400 + 6 负控 400", "12/12 通过"),
        ("公开功能", "6/6 normalized exact", "6/6 精确"),
        ("隐藏功能", "≥ 11/12", "proxy 11/12"),
        ("请求成功率", "≥ 99%（目标 100%）", "200/200 = 100%"),
        ("五项 boolean", "全部为真", "全部为真"),
        ("性能 cell", "success 5/5、CV ≤ 10%", "7/7 有效"),
    ]
    cw = [Inches(1.9), Inches(2.5), Inches(2.0)]
    table(s, Inches(0.7), y, Inches(6.4), rows, cw, row_h=Inches(0.36))

    tf = textbox(s, Inches(7.5), y, Inches(5.1), Inches(0.3))
    para(tf, "关键缺陷：客户端断连导致容量泄漏", size=13, bold=True, color=WARN, first=True)
    bullets(s, Inches(7.5), y + Inches(0.4), Inches(5.0), [
        ("现象 ", "客户端在 8K 请求中断连后，服务连续 653 s 返回 503，"
                 "而 /health 始终 200"),
        ("根因 ", "prefill 完成前无响应字节可写，HTTP 层无法从写失败感知断连，"
                 "取消信号无处生效"),
        ("修复 ", "socket EOF 监视 + prefill 块边界检查取消，"
                 "中止映射为 Cancelled 而非服务故障"),
        ("结果 ", "恢复时间 653 s → 4.86 s；新增 3 个回归测试，"
                 "修复前失败、修复后通过"),
    ], size=12, gap=11)
    footer(s, "无平台批准的 trajectory reference，故 trajectory 标记 unverified，且未使用自捕获自评分")

    # ───────────────────────── 5. Optimization path ─────────────────────────
    s = blank(prs)
    y = header(s, "第二层：性能工程的实际路径", "PERFORMANCE")
    tf = textbox(s, Inches(0.7), y, Inches(11.9), Inches(0.42))
    para(tf, "十七次单变量配对实验，九项接受、八项拒绝。每项接受的改动都带一个 =0 回滚开关，"
             "并以位相等或有界漂移的测试守护。",
         size=13.5, color=INK, first=True, line=1.3)

    y += Inches(0.6)
    rows = [
        ("接受的改动", "定位到的问题", "端到端效果"),
        ("release 构建", "host 侧解释开销", "TPOT −18.3%"),
        ("关闭调试采集", "每 token 写盘与 stderr", "TPOT −3.8%"),
        ("W4 prefill 改 dequant+GEMM", "每行重读整个权重矩阵", "TTFT −74.3%"),
        ("packed W4 GEMV", "128 B 事务只用 16 B", "TPOT −35.4%"),
        ("dequant scratch 池化", "178 MB cudaMalloc 耗 3.2 ms", "TTFT 1K −65%"),
        ("prefill 块 64 → 512", "每块重复 dequant", "TTFT 1K −69%"),
        ("attention 批量化 + 行协作 softmax", "每线程重复扫描整行", "16K TTFT −92%"),
    ]
    cw = [Inches(3.5), Inches(2.5), Inches(1.6)]
    table(s, Inches(0.7), y, Inches(7.6), rows, cw, row_h=Inches(0.355), body_size=11.5)

    tf = textbox(s, Inches(8.6), y, Inches(4.0), Inches(0.3))
    para(tf, "被拒绝的八项（负结果）", size=13, bold=True, color=WARN, first=True)
    bullets(s, Inches(8.6), y + Inches(0.4), Inches(3.9), [
        "deferred GDN 状态检查：零收益",
        "warp-per-output GEMV：−1.8%（噪声内）",
        "uint4 向量化加载：慢 26%",
        "4 累加器数组：慢 17%（local memory 溢出）",
        "展开 + FMA 融合：慢 1.3%",
        "Marlin 位技巧 SIMD：慢 54%",
        "per-group 查表：设计阶段否决",
        "chunk 512 初判保守：显存复核后采纳",
    ], size=11.5, gap=8, bullet="·")
    footer(s, "接受条件：correctness 不降、CV ≤ 10%、端到端收益超过噪声；只有 kernel 变快而端到端无收益一律拒绝")

    # ───────────────────────── 6. Finding: softmax ─────────────────────────
    s = blank(prs)
    y = header(s, "核心发现：16K 的瓶颈是 softmax 的 512 倍冗余扫描", "DEEP DIVE 1")

    rect(s, Inches(0.7), y, Inches(11.9), Inches(1.02), fill=PANEL)
    tf = textbox(s, Inches(0.95), y + Inches(0.15), Inches(11.4), Inches(0.75))
    para(tf, "for (c = 0; c < valid_cols; c++)  // 这一行被 block 内每个线程各执行一遍",
         size=14, bold=True, color=WARN, font=MONO, first=True, space_after=5)
    para(tf, "256 个线程做完全相同的全行扫描，再乘以每行 ceil(cols/256) 个 block —— "
             "每行被完整扫描 512 次，而正确做法只需两遍分区扫描。",
         size=12.5, color=MUTED, space_after=0)

    y += Inches(1.22)
    tf = textbox(s, Inches(0.7), y, Inches(5.9), Inches(0.3))
    para(tf, "定位过程（三次假设修正）", size=13, bold=True, color=ACCENT, first=True)
    bullets(s, Inches(0.7), y + Inches(0.38), Inches(5.7), [
        ("① 单层模型仅解释 124 s / 408.7 s ", "→ 改用逐块递增 kv 的阶梯测量，"
         "单个 attention 层 25.05 s，16 层约 400 s"),
        ("② 假设 GEMM 形状退化 ", "→ 重排后 M 从 6 变 512，仅降至 23.59 s，非主因"),
        ("③ 假设输出 stride 问题 ", "→ 微基准显示 GEMM 仅 5.7 µs，与 stride 无关"),
        ("④ 定位测量口径错误 ", "→ 此前用 kv_offset=0 测 softmax（平均仅 256 列），"
         "按真实 15872 列换算得 2234 ms，与实测 2082 ms 吻合"),
    ], size=12, gap=9)

    tf = textbox(s, Inches(6.9), y, Inches(5.7), Inches(0.3))
    para(tf, "修复效果（双卡双副本验证）", size=13, bold=True, color=GOOD, first=True)
    rows = [
        ("测量项", "修复前", "修复后", "倍数"),
        ("attention 阶梯（单层）", "23.06 s", "0.659 s", "35.0×"),
        ("单次 sdpa（kv=16384）", "2053.7 ms", "13.7 ms", "150×"),
        ("服务 TTFT 8K", "64.1 s", "17.8 s", "3.6×"),
        ("服务 TTFT 16K", "408.7 s", "31.9 s", "12.8×"),
    ]
    cw = [Inches(2.25), Inches(1.15), Inches(1.15), Inches(0.95)]
    table(s, Inches(6.9), y + Inches(0.4), Inches(5.5), rows, cw,
          row_h=Inches(0.36), body_size=11.5)
    tf = textbox(s, Inches(6.9), y + Inches(2.28), Inches(5.6), Inches(0.7))
    para(tf, "改为一个 block 负责一行、256 线程分工 + shared memory 归约；"
             "每个线程只重写自己读过的元素，因而原地安全。max 归约与顺序无关，"
             "指数求和改为分区树形，属于已接受的同类 FP32 重结合。",
         size=11.5, color=MUTED, first=True, line=1.25)
    footer(s, "公开功能 6/6 在修复后重新验证仍全部精确通过")

    # ───────────────────────── 7. Finding: memory & W4 ─────────────────────────
    s = blank(prs)
    y = header(s, "核心发现：两个“非计算”瓶颈", "DEEP DIVE 2")

    tf = textbox(s, Inches(0.7), y, Inches(5.8), Inches(0.32))
    para(tf, "① prefill 的 W4 投影退化成 M 个 GEMV", size=14, bold=True, color=ACCENT, first=True)
    rect(s, Inches(0.7), y + Inches(0.42), Inches(5.7), Inches(0.72), fill=PANEL)
    tf = textbox(s, Inches(0.92), y + Inches(0.54), Inches(5.3), Inches(0.55))
    para(tf, "blockIdx.x = row * out_features + out", size=12.5, bold=True,
         color=WARN, font=MONO, first=True, space_after=3)
    para(tf, "每个 block 各自重读所需权重行 ⇒ M 行 prefill 把整个权重矩阵搬运 M 次",
         size=11.5, color=MUTED, space_after=0)
    bullets(s, Inches(0.7), y + Inches(1.3), Inches(5.7), [
        ("证据 ", "prefill 每 token 76 ms，decode 105 ms；64-token block 只比单步 decode 慢 46 倍"
                 "却做了 64 倍的工作，批处理效率仅 1.03 倍"),
        ("修复 ", "大 M 时先 dequant 到 BF16 scratch 再交 tensor core GEMM，"
                 "decode 保持 GEMV；解压逐位精确（assert_eq! 全量断言）"),
        ("效果 ", "1K TTFT 77.33 s → 19.87 s，公开功能仍 6/6"),
    ], size=12, gap=8)

    tf = textbox(s, Inches(6.9), y, Inches(5.7), Inches(0.32))
    para(tf, "② 80% 的 prefill 花在 cudaMalloc", size=14, bold=True, color=ACCENT, first=True)
    rows = [
        ("微基准（单层归因）", "耗时"),
        ("GDN 层整体 prefill", "16.7 ms"),
        ("  其中 norm + out_proj + MLP", "11.4 ms"),
        ("cudaMalloc+free 178 MB", "3.22 ms"),
        ("cudaMalloc+free 4 B", "3.7 µs"),
        ("单个 MLP 投影（含分配）", "3.92 ms"),
    ]
    cw = [Inches(3.9), Inches(1.5)]
    table(s, Inches(6.9), y + Inches(0.42), Inches(5.4), rows, cw,
          row_h=Inches(0.345), body_size=11.5)
    tf = textbox(s, Inches(6.9), y + Inches(2.55), Inches(5.6), Inches(1.1))
    para(tf, "178 MB 的 dequant scratch 分配+释放要 3.22 ms，占一次完整 MLP 投影的 82%。"
             "每个 GDN 层每块有 7 次 W4 投影，各自新分配一次 —— 约 80% 的 prefill 时间"
             "是显存页表操作，不是算术。", size=12, color=MUTED, first=True, line=1.25)
    para(tf, "改为按尺寸类复用的 scratch 池后：GDN 层 prefill 16.7 → 5.6 ms，"
             "1K TTFT 5.76 s → 2.38 s，且 frozen 输出逐位不变。",
         size=12, bold=True, color=GOOD, space_after=0, line=1.25)
    footer(s, "两处根因都不是“kernel 算得慢”，而是访存与内存管理的结构问题 —— 靠实测归因才暴露")

    # ───────────────────────── 8. Methodology ─────────────────────────
    s = blank(prs)
    y = header(s, "实验方法与结论可信度保障", "METHOD")
    items = [
        ("四卡并行实验",
         "GPU1 只跑正式服务 A/B（带全局 flock），GPU0 做 profile 归因，GPU2/3 跑测试分片。"
         "一次验证周期从串行 35 s 降到 12 s；三种 chunk 配置的扫描一次 11 s 完成。"),
        ("微基准矩阵做成本归因",
         "为 W4 GEMV 构造三个诊断变体（只读权重 / 去 dequant / 完整），一次并行得到"
         "「访存 41% + 乘加 21% + dequant 算术 39%」的分解，并证明访存已达峰值 81%——"
         "瓶颈定位由测量给出，不依赖推测。"),
        ("先证明位相等，再写 kernel",
         "Marlin 位技巧动工前先穷举验证：0x4300|q 恰为 128+q、BF16 减法精确、"
         "乘积 ≤13 位有效可被 FP32 精确保存 ⇒ 与生产路径必然位相等。"
         "正确性由数学保证而非事后调容差。"),
        ("对低于 5% 的结论强制双卡双副本复测",
         "一次三卡测量曾显示 4% 收益，双副本复测得 baseline 923/901 µs、候选 934/913 µs —— "
         "卡间差异 2.4% 大于候选效应，实际慢 1.3%。该协议排除了一次假阳性。"),
    ]
    yy = y
    for i, (t, d) in enumerate(items):
        rect(s, Inches(0.7), yy, Inches(11.9), Inches(1.12), fill=WHITE, line=RULE)
        tf = textbox(s, Inches(1.0), yy + Inches(0.14), Inches(11.3), Inches(0.3))
        para(tf, f"{i + 1}. {t}", size=14, bold=True, color=ACCENT, first=True)
        tf = textbox(s, Inches(1.0), yy + Inches(0.46), Inches(11.3), Inches(0.6))
        para(tf, d, size=12, color=MUTED, first=True, line=1.28)
        yy += Inches(1.22)
    footer(s, "开发期共有三次假设被实测推翻，均记录于 REPORT.md，作为后续决策依据")

    # ───────────────────────── 9. Trade-offs ─────────────────────────
    s = blank(prs)
    y = header(s, "取舍：correctness / 性能 / 稳定性 / 显存", "TRADE-OFFS")
    pairs = [
        ("轨迹软目标 ⇄ TTFT", GOOD,
         "W4 GEMM 与 packed GEMV 改变累加顺序，冻结轨迹前缀在 28 / 23 / 76 token 间变动。"
         "机制唯一：oracle 自身在第 23、28、76 步的 top1/top2 margin 恰为 0，"
         "任何容差都无法修。轨迹阈值为 0 且公开功能始终 6/6，故以 5 分软目标换 35 分主轴。"),
        ("NaN 防护 ⇄ 速度", ACCENT,
         "曾把逐 op 有限性检查批量化为每 token 一次（去掉约 200 次同步）。实测零收益，"
         "于是保留更严格的逐 op 检查 —— 更安全的选项恰好也是免费的。"),
        ("显存 ⇄ prefill 吞吐", WARN,
         "prefill 块 64 → 512 使 attention score workspace 从 192 MB 涨到 1536 MB，"
         "scratch 池另占约 460 MB。两者都经 request_state_bytes 计入 admission，"
         "配置过大时启动即 fail closed 而非请求中途 OOM。峰值 19958 / 24564 MiB。"),
        ("稳定性：不做交换", ACCENT,
         "每项改动接受前都需通过协议 12/12、公开功能 6/6 与干净 soak；"
         "两个在隔离测试中更快的候选因端到端无收益而被拒。没有任何改动只凭 kernel 数字被接受。"),
    ]
    yy = y
    for title, color, desc in pairs:
        rect(s, Inches(0.7), yy, Inches(11.9), Inches(1.12), fill=PANEL)
        tf = textbox(s, Inches(1.0), yy + Inches(0.13), Inches(11.3), Inches(0.3))
        para(tf, title, size=13.5, bold=True, color=color, first=True)
        tf = textbox(s, Inches(1.0), yy + Inches(0.44), Inches(11.3), Inches(0.62))
        para(tf, desc, size=11.8, color=MUTED, first=True, line=1.26)
        yy += Inches(1.22)
    footer(s, "显存账本也决定了 bonus 边界：长上下文 65536 需 4923 MiB 而仅余 4606 MiB，131072 则无论如何不可能")

    # ───────────────────────── 10. Limits & honesty ─────────────────────────
    s = blank(prs)
    y = header(s, "已知限制与评测边界", "LIMITS")
    tf = textbox(s, Inches(0.7), y, Inches(11.9), Inches(0.42))
    para(tf, "以下各项均已完整写入 REPORT.md。", size=13.5, color=INK, first=True)

    y += Inches(0.55)
    tf = textbox(s, Inches(0.7), y, Inches(5.8), Inches(0.3))
    para(tf, "不作声称的结论", size=13, bold=True, color=WARN, first=True)
    bullets(s, Inches(0.7), y + Inches(0.38), Inches(5.7), [
        ("无正式 scorer 产物 ", "评分需平台批准的 trajectory reference，"
                              "本机不存在；未以自捕获方式自评分"),
        ("因此不声称 ", "eligible、trajectory 得分、基础分或多模态得分"),
        ("no_xid 为间接证据 ", "dmesg 无读取权限，以服务日志零 CUDA 错误 + "
                            "GPU 计数稳定佐证，已标注为间接"),
        ("hidden 用自建代理集 ", "官方 hidden 集不可得，11/12 来自独立种子代理集"),
        ("图片套件不在本机 ", "public 4/4 属平台运行项；自建探针非官方公开集"),
    ], size=12, gap=10)

    tf = textbox(s, Inches(6.9), y, Inches(5.8), Inches(0.3))
    para(tf, "性能与 bonus 的边界", size=13, bold=True, color=WARN, first=True)
    bullets(s, Inches(6.9), y + Inches(0.38), Inches(5.7), [
        ("与合同标尺的距离 ", "prefill 距算术下限约 4.8 倍，decode 距带宽下限"
                            "约 3.3 倍；动态分按同轮最优参考计分"),
        ("decode 达局部最优 ", "五次微优化实测均无效，反汇编证明 M=1 GEMV "
                             "受延迟限制；突破需 M≥8 的融合 dequant-MMA"),
        ("长上下文受显存限制 ", "65536 需 4923 MiB、可用 4606 MiB；131072 不可行"),
        ("C8 与 MTP 未交付 ", "C4 已交付；C8 需 4.1 GiB 超出预算，"
                            "共享 scratch 的改造路径已量化但未实现"),
    ], size=12, gap=10)
    footer(s, "八项被拒实验均保留数据与根因分析，作为负结果证据记录于 REPORT.md")

    # ───────────────────────── 11. Multimodal delivered ─────────────────────────
    s = blank(prs)
    y = header(s, "多模态 bonus：完整交付", "BONUS")
    tf = textbox(s, Inches(0.7), y, Inches(11.9), Inches(0.42))
    para(tf, "完整 vision 路径 + POST /v1/chat/completions 已接入服务；"
             "APXINF_ENABLE_MULTIMODAL=1 单开关启用，默认文本配置行为与数值零变化。",
         size=13.5, color=INK, first=True)

    y += Inches(0.56)
    tf = textbox(s, Inches(0.7), y, Inches(5.8), Inches(0.3))
    para(tf, "正确性验证（逐阶段）", size=13, bold=True, color=GOOD, first=True)
    bullets(s, Inches(0.7), y + Inches(0.38), Inches(5.7), [
        ("预处理位相等 ", "PNG 解码 + patchify 与 HF processor 逐位一致"
                        "（1,204,224 个 f32 零差异）"),
        ("prompt token 级精确 ", "模板渲染 + tokenize 与 Qwen3VLProcessor "
                              "在全部探针上一致"),
        ("vision 塔对齐 oracle ", "27 块 BF16 对 FP32 golden 漂移 1.2–8%，"
                               "merged 嵌入余弦 0.99972"),
        ("新 kernel 位相等回归 ", "宽 head_dim SDPA 在 64 处与旧 kernel 逐位一致；"
                               "merger 采用 erf-GELU，经 HF 源码核实"),
    ], size=12, gap=10)

    tf = textbox(s, Inches(6.9), y, Inches(5.8), Inches(0.3))
    para(tf, "端到端实测结果", size=13, bold=True, color=ACCENT, first=True)
    bullets(s, Inches(6.9), y + Inches(0.38), Inches(5.7), [
        ("自建图片探针 4/4 精确 ", "\"382\" / \"red\" / \"12\" / \"5\"，"
                               "单请求端到端 1–2 s"),
        ("文本门禁复验通过 ", "协议 12/12、公开功能 6/6、图文混合 soak 36/36"),
        ("文本数值零变化 ", "冻结 128-token 输出与文本配置逐字节一致"),
        ("显存过账 ", "vision 塔约 956 MiB 常驻，chunk 自动降为 256，"
                    "峰值 20914 / 24564 MiB"),
    ], size=12, gap=10)

    rect(s, Inches(0.7), Inches(5.62), Inches(11.9), Inches(1.1), fill=PANEL)
    tf = textbox(s, Inches(0.95), Inches(5.76), Inches(11.4), Inches(0.85))
    para(tf, "得分口径说明", size=12.5, bold=True, color=WARN, first=True, space_after=4)
    para(tf, "官方图片套件不在本机，「public 4/4」属平台运行项；自建探针与合同同构同格式"
             "但非官方公开集，因此不声称多模态得分。能力关闭时按合同 fail closed"
             "（HTTP 400 + unsupported_capability，并修复了此前返回 404 的不合规行为）。",
         size=11.5, color=MUTED, space_after=0, line=1.3)
    footer(s, "回滚 = 取消 APXINF_ENABLE_MULTIMODAL 单开关；文本 layer-1/layer-2 全部性能数字均在无该开关的配置下测得，不受影响")

    # ───────────────────────── 11.5 C4 concurrent bonus ─────────────────────────
    s = blank(prs)
    y = header(s, "多请求 bonus：C4 四路并发交付", "BONUS")
    tf = textbox(s, Inches(0.7), y, Inches(11.9), Inches(0.42))
    para(tf, "批式 protocol runtime 仅在 APXINF_Q35_MAX_CONCURRENCY>1 时启用，默认单请求路径零改动；"
             "官方 evaluator 校准 run 全部效度门通过。",
         size=13.5, color=INK, first=True)

    y += Inches(0.56)
    tf = textbox(s, Inches(0.7), y, Inches(5.8), Inches(0.3))
    para(tf, "实现要点", size=13, bold=True, color=GOOD, first=True)
    bullets(s, Inches(0.7), y + Inches(0.38), Inches(5.7), [
        ("批式调度 ", "admission 4 permits + 单 batch worker，每轮一步批式 decode；"
                    "两阶段交付消除 permit 释放竞态"),
        ("批式 kernel 位相等 ", "GDN conv/recurrent、partial RoPE 批式版与串行路径"
                             "逐字节一致（真实权重回归断言）"),
        ("prefix cache + 会话回收池 ", "相同 prompt 从模板 fork 约 2 ms；"
                                    "流序分配把冷 fork 从 330 ms 降到个位数 ms"),
        ("单请求零回归 ", "TPOT 在冻结基线带内，输出字节不变；"
                       "batch=1 自动回退串行路径"),
    ], size=12, gap=10)

    tf = textbox(s, Inches(6.9), y, Inches(5.8), Inches(0.3))
    para(tf, "官方校准结果（multi-c4-text-perf-1024）", size=13, bold=True, color=ACCENT, first=True)
    rows = [
        ("指标", "实测", "门槛"),
        ("success / correctness", "1.0 / 1.0", "= 1.0"),
        ("Jain 公平性", "0.9993", ">= 0.95"),
        ("p95 TTFT", "1.789 s", "<= 2.496 s"),
        ("p95 TPOT", "173.9 ms", "<= 186.9 ms"),
        ("goodput", "23.53 tok/s", "计分指标"),
        ("no_fallback / 健康", "true", "必须为真"),
    ]
    table(s, Inches(6.9), y + Inches(0.4), Inches(5.7),
          rows, [Inches(2.3), Inches(1.7), Inches(1.7)], row_h=Inches(0.34))

    rect(s, Inches(0.7), Inches(5.62), Inches(11.9), Inches(1.1), fill=PANEL)
    tf = textbox(s, Inches(0.95), Inches(5.76), Inches(11.4), Inches(0.85))
    para(tf, "C8 未交付原因", size=12.5, bold=True, color=WARN, first=True, space_after=4)
    para(tf, "受限于显存而非调度：8 份完整请求态需 4.1 GiB，超出 2.36 GiB 的 admission 预算；"
             "共享 GDN scratch 后每份可降至约 295 MiB、可满足入场条件，但该改造未实现，"
             "因此不声称 C8。调度器、admission 与批式 kernel 均已按并发数参数化。",
         size=11.5, color=MUTED, space_after=0, line=1.3)
    footer(s, "回滚 = 不设 APXINF_Q35_MAX_CONCURRENCY（或 =1）即恢复精确的冻结单请求 runtime；"
              "并发同批 32 请求全部 200，无 503")

    # ───────────────────────── 12. Reproduce ─────────────────────────
    s = blank(prs)
    y = header(s, "可复现性与提交材料", "REPRODUCIBILITY")

    rect(s, Inches(0.7), y, Inches(11.9), Inches(2.02), fill=RGBColor(0x1A, 0x1D, 0x24))
    tf = textbox(s, Inches(1.0), y + Inches(0.14), Inches(11.3), Inches(1.78))
    for i, line in enumerate([
        "# 提交 commit（验收以此为准）",
        "06993a2d2642c6f7177b57493b797d5d537e4d64    branch: integrate/member2",
        "",
        "python3 benchmarks/qwen38_4090/evaluation/test.py check      # → assignment checks passed",
        "cargo build --release --features cuda-no-nvtx --locked --bin apxinf",
        "target/release/apxinf serve --model <ckpt> --revision 63768c10... --gpu-uuid GPU-343bc895...",
        "APXINF_ENABLE_MULTIMODAL=1 target/release/apxinf serve ...   # 多模态配置（可选，单开关）",
        "APXINF_Q35_MAX_CONCURRENCY=4 ... serve ... --queue-capacity 4  # C4 并发配置（可选，单开关）",
    ]):
        color = RGBColor(0x7E, 0xC8, 0xFF) if line.startswith("#") else RGBColor(0xE6, 0xEA, 0xF0)
        para(tf, line, size=11.5, color=color, font=MONO, first=(i == 0), space_after=2)

    y += Inches(2.2)
    tf = textbox(s, Inches(0.7), y, Inches(5.8), Inches(0.3))
    para(tf, "提交材料对应关系", size=13, bold=True, color=ACCENT, first=True)
    bullets(s, Inches(0.7), y + Inches(0.38), Inches(5.7), [
        ("设计变化与执行阶段 ", "REPORT.md 按第一层/第二层分节，逐项记录"),
        ("test.py check / run ", "check 已通过并留存日志；run 因缺平台参照件无法产出评分件，等价负载已直测"),
        ("负控制与回归测试 ", "7 个协议负控、3 个断连故障注入、9 组数值回归"),
        ("取舍与限制 ", "REPORT.md 专设 Trade-offs 与 Limits 两节"),
    ], size=12, gap=9)

    tf = textbox(s, Inches(6.9), y, Inches(5.8), Inches(0.3))
    para(tf, "回滚能力", size=13, bold=True, color=ACCENT, first=True)
    bullets(s, Inches(6.9), y + Inches(0.38), Inches(5.7), [
        ("每项优化独立可关 ", "APXINF_Q35_W4_PREFILL_GEMM / W4_PACKED_GEMV / "
                           "BATCHED_SDPA / SCRATCH_POOL / PREFILL_CHUNK / ROWWISE_SOFTMAX"),
        ("bonus 均单开关 ", "APXINF_ENABLE_MULTIMODAL 与 APXINF_Q35_MAX_CONCURRENCY "
                         "都不设时，行为即纯文本单请求提交件"),
        ("被拒候选保留在树中 ", "默认关闭并附回归测试，作为负结果证据"),
        ("回滚点 ", "47ec280d2f88e8daf87750c0957e596e3a5390c1（集成前 HEAD）"),
        ("清洁性 ", "cargo fmt --check、cargo check --workspace --locked、git diff --check 全绿"),
    ], size=12, gap=9)
    footer(s, "未提交模型权重、凭据、机器地址或未公开评测数据；submission.json 的图片字段未手工填写")

    # ───────────────────────── 12.5 Division of work ─────────────────────────
    s = blank(prs)
    y = header(s, "分工说明", "TEAM")
    tf = textbox(s, Inches(0.7), y, Inches(11.9), Inches(0.42))
    para(tf, "两人团队，分工依据 git 提交记录。", size=13.5, color=INK, first=True)

    y += Inches(0.52)
    tf = textbox(s, Inches(0.7), y, Inches(7.4), Inches(0.34))
    para(tf, "程仁龙（负责人）", size=14.5, bold=True, color=ACCENT, first=True)
    bullets(s, Inches(0.7), y + Inches(0.46), Inches(7.3), [
        ("模型与 kernel ", "Qwen3.8 config / loader / weights，48 层 GDN 与 "
                         "16 层全注意力 CUDA 实现，W4A16 解包与 GEMM / GEMV kernel"),
        ("服务与协议 ", "HTTP / SSE、七项负控、admission、断连取消与故障恢复，"
                      "/v1/chat/completions 图文入口"),
        ("正确性验证 ", "离线 oracle（文本 + vision）、hidden 代理集、"
                      "逐层对拍与位相等断言"),
        ("性能工程 ", "十七次单变量配对实验、显存账本、四卡并行验证方法"),
        ("双 bonus 与集成 ", "vision 塔 CUDA 移植、C4 批式调度与批式 kernel、"
                          "REPORT.md、最终集成与提交"),
    ], size=12.5, gap=11)

    tf = textbox(s, Inches(8.6), y, Inches(4.0), Inches(0.34))
    para(tf, "王天民", size=14.5, bold=True, color=ACCENT, first=True)
    bullets(s, Inches(8.6), y + Inches(0.46), Inches(3.9), [
        ("离线 benchmark 脚手架 ", "shape inventory 与 "
                                "experiment validator 脚本"),
        ("campaign 资料 ", "manifest 与 README"),
        ("协作记录 ", "交接文档"),
    ], size=12.5, gap=11)
    footer(s, "原计划的多人并行安排实际为单人串行推进；GPU 带锁队列与单卡正式重放的纪律全程执行")

    # ───────────────────────── 13. Closing ─────────────────────────
    s = blank(prs)
    rect(s, 0, 0, W, H, fill=RGBColor(0x0F, 0x2A, 0x4A))
    tf = textbox(s, Inches(1.1), Inches(1.35), Inches(11.1), Inches(0.5))
    para(tf, "为什么叫「4090 能飞」", size=15, bold=True,
         color=RGBColor(0x8F, 0xC2, 0xF0), first=True)
    tf = textbox(s, Inches(1.1), Inches(1.88), Inches(11.1), Inches(2.1))
    para(tf, "让它飞起来的不是更强的硬件，是看清了瓶颈的真实位置。",
         size=22, bold=True, color=WHITE, first=True, space_after=10, line=1.3)
    para(tf, "三个最大的瓶颈都不在算术里：softmax 每个线程重复扫描整行、"
             "W4 投影每行重读整个权重矩阵、80% 的 prefill 耗在 cudaMalloc。"
             "未更换硬件与算法，通过正确的测量归因逐一消除结构性瓶颈 —— "
             "TTFT 约 47 倍、TPOT 约 2 倍，correctness 与稳定性全程未退让。",
         size=15, color=RGBColor(0xC5, 0xD9, 0xEE), space_after=0, line=1.35)

    yy = Inches(4.5)
    for i, (k, v) in enumerate([
        ("最大的单点收益", "行协作 softmax：16K attention 单层 23.06 s → 0.659 s"),
        ("最有价值的负结果", "Marlin 位技巧在 M=1 GEMV 慢 54%，反汇编级根因"),
        ("最关键的方法", "低于 5% 的结论必须双卡双副本复测"),
    ]):
        rect(s, Inches(1.1) + Inches(3.75) * i, yy, Inches(3.55), Inches(1.5),
             fill=RGBColor(0x16, 0x38, 0x60))
        tf = textbox(s, Inches(1.32) + Inches(3.75) * i, yy + Inches(0.2),
                     Inches(3.15), Inches(0.3))
        para(tf, k, size=12, bold=True, color=RGBColor(0x8F, 0xC2, 0xF0), first=True)
        tf = textbox(s, Inches(1.32) + Inches(3.75) * i, yy + Inches(0.56),
                     Inches(3.15), Inches(0.85))
        para(tf, v, size=12, color=WHITE, first=True, line=1.3)

    tf = textbox(s, Inches(1.1), Inches(6.5), Inches(11.1), Inches(0.4))
    para(tf, "commit 06993a2d2642c6f7177b57493b797d5d537e4d64 · branch integrate/member2",
         size=11, color=RGBColor(0x8F, 0xC2, 0xF0), font=MONO, first=True)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"wrote {OUT}")
    print(f"slides: {len(prs.slides.__iter__.__self__._sldIdLst)}")


if __name__ == "__main__":
    build()
